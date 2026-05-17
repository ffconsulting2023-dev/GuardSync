#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
GuardSync プロジェクト計画書 v1.1 更新スクリプト
MAXプラン追加要件（SaaS販売・FC展開・契約違反監視・人員融通優先）を反映
"""

import os

OUTPUT_DIR = "/Users/t.take/AI management/project/output/GuardSync"

MAX_NEW_FEATURES = [
    {
        "category": "SaaS販売機能",
        "summary": "他社警備会社へのSaaS販売・契約・課金管理",
        "details": [
            "料金プラン設定（ライセンス数・機能制限・価格）",
            "契約期間・自動更新管理",
            "Stripe連携（クレジットカード決済・自動請求）",
            "利用状況レポート（隊員数・ポスト数・API使用量）",
            "請求書自動発行・入金管理",
        ],
        "priority": "v2",
        "risk": "Stripe連携の設定・テストに時間が必要"
    },
    {
        "category": "独立支援・FC展開管理",
        "summary": "フランチャイズ本部・加盟店の階層管理と独立支援制度",
        "details": [
            "FC本部・加盟店の階層管理（親子テナント構造）",
            "独立支援パッケージ（初期設定・研修資料の配布）",
            "FC契約書の電子締結・管理（電子契約機能と連携）",
            "ロイヤリティ計算・徴収管理（売上連動型・固定型）",
            "加盟店への本部指導記録・お知らせ一括配信",
        ],
        "priority": "v2",
        "risk": "FC契約の法的要件確認が必要（フランチャイズ法）"
    },
    {
        "category": "契約違反監視（スーパー管理者）",
        "summary": "FC・SaaS加盟店の契約遵守状況を開発側が監視・データ抽出",
        "details": [
            "全テナントKPI監視ダッシュボード（売上・隊員数・稼働率）",
            "契約違反検知アラート（規約禁止行為・異常データパターン）",
            "FC契約条件の逸脱自動検知（売上閾値・人員数制限等）",
            "テナント別データのCSV/Excel抽出（証拠保全用）",
            "抽出操作の監査ログ（誰がいつ何を抽出したか）",
            "スーパー管理者限定アクセス・操作ログ必須",
        ],
        "priority": "v2",
        "risk": "利用規約・プライバシーポリシーにデータアクセス同意条項が必要（法的必須）"
    },
    {
        "category": "協力企業・グループ企業 優先人員融通",
        "summary": "企業間の優先順位に基づく自動人員融通打診・実績管理",
        "details": [
            "企業関係性の定義（グループ企業 > 協力企業 > 一般協力会社）",
            "人員不足時の優先順位自動打診フロー",
            "回答期限設定・未回答時の次順位エスカレーション",
            "融通実績管理（企業別・回数・人数・日付）",
            "相互融通バランスの可視化（貸し借り）",
            "融通条件設定（企業間単価・優先/除外現場）",
        ],
        "priority": "v2",
        "risk": "企業間の契約・単価設定が複雑になる可能性"
    },
]


def update_word():
    from docx import Document
    from docx.shared import Pt, RGBColor, Cm
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    path_old = os.path.join(OUTPUT_DIR, "20260518_GuardSync_プロジェクト計画書.docx")
    doc = Document(path_old)

    def add_heading(text, level=1, color=None):
        h = doc.add_heading(text, level=level)
        if color:
            for run in h.runs:
                run.font.color.rgb = RGBColor(*color)
        return h

    def add_table(headers, rows):
        table = doc.add_table(rows=1, cols=len(headers))
        table.style = 'Table Grid'
        hdr = table.rows[0].cells
        for i, h in enumerate(headers):
            hdr[i].text = h
            hdr[i].paragraphs[0].runs[0].font.bold = True
            hdr[i].paragraphs[0].runs[0].font.size = Pt(9)
            hdr[i].paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER
        for row_data in rows:
            row = table.add_row().cells
            for i, val in enumerate(row_data):
                row[i].text = str(val)
                row[i].paragraphs[0].runs[0].font.size = Pt(9)
        return table

    doc.add_page_break()

    add_heading("付録：MAXプラン 追加要件（v1.1更新）", 1, (0x1A, 0x56, 0xDB))
    doc.add_paragraph(
        "2026年5月18日追加。他社への販売・FC展開・契約違反監視・優先人員融通の4機能をv2（8月末）に追加する。"
    )

    for feat in MAX_NEW_FEATURES:
        add_heading(feat["category"], 2)
        doc.add_paragraph(feat["summary"])

        p = doc.add_paragraph()
        p.add_run("主な機能：").font.bold = True
        for detail in feat["details"]:
            bp = doc.add_paragraph(style='List Bullet')
            bp.add_run(detail).font.size = Pt(9)

        note_p = doc.add_paragraph()
        note_run = note_p.add_run(f"優先度：{feat['priority']}　|　リスク：{feat['risk']}")
        note_run.font.size = Pt(9)
        note_run.font.color.rgb = RGBColor(0xDC, 0x26, 0x26)

    add_heading("MAXプラン 全体構成（更新版）", 2)
    rows = [
        ["SaaS販売機能", "v2", "プラン管理・Stripe決済・利用レポート・請求自動発行"],
        ["独立支援・FC展開", "v2", "本部/加盟店管理・ロイヤリティ・FC契約電子締結"],
        ["契約違反監視", "v2", "スーパー管理者ダッシュボード・データ抽出・監査ログ"],
        ["優先人員融通", "v2", "グループ/協力企業の優先打診・融通実績・バランス管理"],
        ["給与集計・労務管理", "v2", "変形労働制・残業・有給・就労時間アラート"],
        ["電子契約完全版", "v2", "紙OCR取込・決裁フロー・タイムスタンプ・人員融通契約"],
        ["帳票強化", "v2", "Excel出力・会計ソフトCSV（マネーフォワード/弥生/freee）"],
        ["隊員アプリ強化", "v2", "休暇申請・出勤希望・周辺施設検索"],
        ["AI自動配置", "v3", "過去データ学習による自動シフト作成"],
        ["会計ソフトAPI連携", "v3", "マネーフォワード・freee・弥生へのAPI直接連携"],
        ["IVR（自動音声電話）", "v3", "ガラケー対応の上下番自動音声報告"],
        ["年末調整", "v3", "年末調整書類の電子管理・計算"],
    ]
    add_table(["機能カテゴリ", "バージョン", "主な機能内容"], rows)

    path_new = os.path.join(OUTPUT_DIR, "20260518_GuardSync_プロジェクト計画書.docx")
    doc.save(path_new)
    print(f"Word更新完了: {path_new}")


def update_excel():
    from openpyxl import load_workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter

    path = os.path.join(OUTPUT_DIR, "20260518_GuardSync_プロジェクト計画書.xlsx")
    wb = load_workbook(path)

    NAVY = "1A56DB"
    LIGHT_BLUE = "DBEAFE"
    LIGHT_GREEN = "DCFCE7"
    LIGHT_RED = "FEE2E2"
    LIGHT_YELLOW = "FEF9C3"
    PURPLE = "F3E8FF"
    WHITE = "FFFFFF"
    GRAY = "6B7280"

    header_font = Font(bold=True, color="FFFFFF", size=10)
    header_fill = PatternFill("solid", fgColor=NAVY)
    body_font = Font(size=9)
    center_align = Alignment(horizontal="center", vertical="center", wrap_text=True)
    wrap_align = Alignment(vertical="center", wrap_text=True)
    thin_border = Border(
        left=Side(style="thin", color="CCCCCC"),
        right=Side(style="thin", color="CCCCCC"),
        top=Side(style="thin", color="CCCCCC"),
        bottom=Side(style="thin", color="CCCCCC"),
    )

    # ====== 新シート: MAXプラン追加要件 ======
    ws_max = wb.create_sheet("MAXプラン追加要件")

    # タイトル
    ws_max.merge_cells("A1:G1")
    ws_max["A1"] = "MAXプラン 追加要件（v1.1 / 2026-05-18追加）"
    ws_max["A1"].font = Font(bold=True, size=14, color="FFFFFF")
    ws_max["A1"].fill = PatternFill("solid", fgColor="7C3AED")
    ws_max["A1"].alignment = Alignment(horizontal="center", vertical="center")
    ws_max.row_dimensions[1].height = 35

    headers = ["カテゴリ", "概要", "主な機能", "バージョン", "リスク・注意事項"]
    for col, h in enumerate(headers, 1):
        c = ws_max.cell(row=2, column=col, value=h)
        c.font = header_font
        c.fill = header_fill
        c.alignment = center_align
        c.border = thin_border

    cat_colors = {
        "SaaS販売機能": LIGHT_BLUE,
        "独立支援・FC展開管理": LIGHT_GREEN,
        "契約違反監視（スーパー管理者）": LIGHT_RED,
        "協力企業・グループ企業 優先人員融通": PURPLE,
    }

    row = 3
    for feat in MAX_NEW_FEATURES:
        details_str = "\n".join([f"・{d}" for d in feat["details"]])
        data = [
            feat["category"],
            feat["summary"],
            details_str,
            feat["priority"],
            feat["risk"],
        ]
        fill_color = cat_colors.get(feat["category"], WHITE)
        for col, val in enumerate(data, 1):
            c = ws_max.cell(row=row, column=col, value=val)
            c.font = body_font
            c.fill = PatternFill("solid", fgColor=fill_color)
            c.border = thin_border
            c.alignment = center_align if col in [4] else wrap_align
        ws_max.row_dimensions[row].height = 80
        row += 1

    col_widths = [25, 35, 50, 10, 40]
    for col, w in enumerate(col_widths, 1):
        ws_max.column_dimensions[get_column_letter(col)].width = w

    # ====== 新シート: MAXプラン全体構成 ======
    ws_plan = wb.create_sheet("MAXプラン全体構成")

    ws_plan.merge_cells("A1:F1")
    ws_plan["A1"] = "MAXプラン 全体機能構成（v1.1更新版）"
    ws_plan["A1"].font = Font(bold=True, size=14, color="FFFFFF")
    ws_plan["A1"].fill = PatternFill("solid", fgColor=NAVY)
    ws_plan["A1"].alignment = Alignment(horizontal="center", vertical="center")
    ws_plan.row_dimensions[1].height = 35

    plan_headers = ["#", "機能カテゴリ", "バージョン", "主な機能内容", "備考"]
    for col, h in enumerate(plan_headers, 1):
        c = ws_plan.cell(row=2, column=col, value=h)
        c.font = header_font
        c.fill = header_fill
        c.alignment = center_align
        c.border = thin_border

    all_features = [
        (1, "SaaS販売機能", "v2", "プラン管理・Stripe決済・利用レポート・請求自動発行", "🆕 追加"),
        (2, "独立支援・FC展開管理", "v2", "本部/加盟店管理・ロイヤリティ計算・FC契約電子締結", "🆕 追加"),
        (3, "契約違反監視（スーパー管理者）", "v2", "全テナントKPI監視・違反検知・データ抽出・監査ログ", "🆕 追加"),
        (4, "優先人員融通管理", "v2", "グループ/協力企業の優先打診・融通実績・バランス管理", "🆕 追加"),
        (5, "給与集計・残業計算", "v2", "日給月給・変形労働制・残業・給与明細メール送信", "既存"),
        (6, "労務管理", "v2", "就労時間アラート・有給管理・有給自動付与", "既存"),
        (7, "電子契約（完全版）", "v2", "紙OCR取込・社内決裁・タイムスタンプ・人員融通契約", "既存"),
        (8, "帳票強化", "v2", "帳票Excel出力・会計ソフトCSV（3社）", "既存"),
        (9, "隊員アプリ強化", "v2", "休暇申請・出勤希望・周辺施設検索", "既存"),
        (10, "現場ごと給与設定", "v2", "現場単位の給与単価管理", "既存"),
        (11, "AI自動配置", "v3", "過去データ学習による自動シフト作成", "既存"),
        (12, "会計ソフトAPI連携", "v3", "マネーフォワード・freee・弥生API直接連携", "既存"),
        (13, "IVR（自動音声電話）", "v3", "ガラケー対応の上下番自動音声報告", "既存"),
        (14, "年末調整", "v3", "年末調整書類の電子管理・計算（社労士確認必須）", "既存"),
        (15, "日払い外部連携", "v3", "即給By GMO・JOBPAY API連携", "既存"),
        (16, "認定タイムスタンプ", "v3", "セイコーソリューションズ等の認定TSA", "既存"),
    ]

    ver_colors = {"v2": LIGHT_GREEN, "v3": LIGHT_YELLOW}
    new_colors = {"🆕 追加": LIGHT_BLUE, "既存": WHITE}

    for i, feat in enumerate(all_features, start=3):
        data = list(feat)
        fill = ver_colors.get(feat[2], WHITE)
        for col, val in enumerate(data, 1):
            c = ws_plan.cell(row=i, column=col, value=val)
            c.font = body_font
            c.fill = PatternFill("solid", fgColor=fill)
            c.border = thin_border
            c.alignment = center_align if col in [1, 3, 5] else wrap_align
        if feat[4] == "🆕 追加":
            for col in range(1, 6):
                ws_plan.cell(row=i, column=col).fill = PatternFill("solid", fgColor="EDE9FE")
        ws_plan.row_dimensions[i].height = 35

    col_widths_plan = [5, 30, 10, 50, 10]
    for col, w in enumerate(col_widths_plan, 1):
        ws_plan.column_dimensions[get_column_letter(col)].width = w

    ws_plan.auto_filter.ref = ws_plan.dimensions

    wb.save(path)
    print(f"Excel更新完了: {path}")


def update_pptx():
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    path = os.path.join(OUTPUT_DIR, "20260518_GuardSync_プロジェクト計画書.pptx")
    prs = Presentation(path)

    blank_layout = prs.slide_layouts[6]

    NAVY = RGBColor(0x1A, 0x56, 0xDB)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT_BLUE = RGBColor(0xDB, 0xEA, 0xFE)
    DARK = RGBColor(0x1F, 0x29, 0x37)
    GRAY = RGBColor(0x6B, 0x72, 0x80)
    PURPLE = RGBColor(0x7C, 0x3A, 0xED)
    LIGHT_PURPLE = RGBColor(0xED, 0xE9, 0xFE)
    GREEN = RGBColor(0x05, 0x96, 0x69)
    LIGHT_GREEN = RGBColor(0xDC, 0xFC, 0xE7)
    RED = RGBColor(0xDC, 0x26, 0x26)
    LIGHT_RED = RGBColor(0xFE, 0xE2, 0xE2)

    def add_rect(slide, x, y, w, h, fill_color, line_color=None):
        shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        if line_color:
            shape.line.color.rgb = line_color
        else:
            shape.line.fill.background()
        return shape

    def add_text_box(slide, text, x, y, w, h, font_size=12, bold=False,
                     color=None, align=PP_ALIGN.LEFT, italic=False):
        txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        if color:
            run.font.color.rgb = color
        return txBox

    def slide_header(slide, title, subtitle=None, color=None):
        hdr_color = color or NAVY
        add_rect(slide, 0, 0, 13.33, 1.2, hdr_color)
        add_text_box(slide, title, 0.4, 0.15, 12, 0.7,
                     font_size=28, bold=True, color=WHITE)
        if subtitle:
            add_text_box(slide, subtitle, 0.4, 0.75, 12, 0.4,
                         font_size=13, color=LIGHT_BLUE)
        add_rect(slide, 0, 7.3, 13.33, 0.2, hdr_color)
        add_text_box(slide, "GuardSync プロジェクト計画書 v1.1 | 2026.05.18",
                     0.4, 7.3, 12, 0.2, font_size=8, color=WHITE,
                     align=PP_ALIGN.CENTER)

    # ====== スライド11: MAXプラン概要 ======
    s11 = prs.slides.add_slide(blank_layout)
    slide_header(s11, "MAXプラン（v2）— 追加要件", "他社販売・FC展開・監視・優先人員融通", PURPLE)

    new_features = [
        ("🏪", "SaaS販売機能", "他社警備会社への販売\nStripe決済・プラン管理・利用レポート", NAVY, LIGHT_BLUE),
        ("🏢", "独立支援・FC展開", "フランチャイズ本部/加盟店管理\nロイヤリティ計算・独立支援パッケージ", GREEN, LIGHT_GREEN),
        ("🔍", "契約違反監視", "スーパー管理者ダッシュボード\n全テナントKPI監視・データ抽出・監査ログ", RED, LIGHT_RED),
        ("🤝", "優先人員融通", "グループ/協力企業を優先した\n自動打診フロー・融通実績・バランス管理", PURPLE, LIGHT_PURPLE),
    ]

    for i, (icon, title, desc, color, bg) in enumerate(new_features):
        x = 0.4 + (i % 2) * 6.4
        y = 1.4 + (i // 2) * 2.85
        add_rect(s11, x, y, 6.1, 2.55, bg)
        add_rect(s11, x, y, 6.1, 0.55, color)
        add_text_box(s11, f"{icon}  {title}", x + 0.15, y + 0.07, 5.8, 0.45,
                     font_size=16, bold=True, color=WHITE)
        add_text_box(s11, desc, x + 0.2, y + 0.7, 5.7, 1.7,
                     font_size=11, color=DARK)

    # ====== スライド12: FC展開詳細 ======
    s12 = prs.slides.add_slide(blank_layout)
    slide_header(s12, "FC展開 + 契約違反監視 詳細", "独立支援制度とスーパー管理者機能", PURPLE)

    add_rect(s12, 0.3, 1.3, 6.1, 5.9, LIGHT_GREEN)
    add_rect(s12, 0.3, 1.3, 6.1, 0.5, GREEN)
    add_text_box(s12, "🏢 独立支援・FC展開管理", 0.45, 1.35, 5.8, 0.4,
                 font_size=14, bold=True, color=WHITE)

    fc_items = [
        "FC本部・加盟店の階層管理（親子テナント）",
        "独立支援パッケージ（初期設定・研修資料）",
        "FC契約書の電子締結・管理",
        "ロイヤリティ計算・徴収管理（売上連動/固定）",
        "加盟店への本部指導記録・一括お知らせ",
    ]
    for j, item in enumerate(fc_items):
        add_text_box(s12, f"✓  {item}", 0.45, 1.95 + j * 0.95, 5.8, 0.8,
                     font_size=10, color=DARK)

    add_rect(s12, 6.9, 1.3, 6.1, 5.9, LIGHT_RED)
    add_rect(s12, 6.9, 1.3, 6.1, 0.5, RED)
    add_text_box(s12, "🔍 契約違反監視（スーパー管理者）", 7.05, 1.35, 5.8, 0.4,
                 font_size=14, bold=True, color=WHITE)

    monitor_items = [
        "全テナントKPI監視ダッシュボード",
        "売上・隊員数・稼働率の推移確認",
        "契約違反・異常データの自動アラート",
        "FC契約条件の逸脱自動検知",
        "テナント別データCSV/Excel抽出",
        "⚠️ 利用規約にデータアクセス同意条項必須",
    ]
    for j, item in enumerate(monitor_items):
        color = RED if "⚠️" in item else DARK
        add_text_box(s12, f"{'⚠️' if '⚠️' in item else '✓'}  {item.replace('⚠️ ', '')}",
                     7.05, 1.95 + j * 0.95, 5.8, 0.8,
                     font_size=10, color=color)

    # ====== スライド13: 人員融通詳細 ======
    s13 = prs.slides.add_slide(blank_layout)
    slide_header(s13, "優先人員融通管理 詳細", "協力企業・グループ企業との連携強化", PURPLE)

    add_rect(s13, 0.3, 1.3, 12.7, 5.9, LIGHT_PURPLE)

    add_text_box(s13, "優先順位の設定", 0.5, 1.4, 4.0, 0.5,
                 font_size=14, bold=True, color=PURPLE)
    priority_items = [
        ("1位", "グループ企業", "最優先・特別単価"),
        ("2位", "協力企業A", "優先・契約単価"),
        ("3位", "協力企業B以下", "通常・一般単価"),
    ]
    for j, (rank, name, note) in enumerate(priority_items):
        y = 2.0 + j * 0.9
        add_rect(s13, 0.4, y, 3.8, 0.75, PURPLE if j == 0 else
                 RGBColor(0xA7, 0x8B, 0xFA) if j == 1 else RGBColor(0xC4, 0xB5, 0xFD))
        add_text_box(s13, rank, 0.5, y + 0.1, 0.8, 0.5,
                     font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(s13, name, 1.3, y + 0.05, 1.8, 0.35,
                     font_size=12, bold=True, color=WHITE)
        add_text_box(s13, note, 1.3, y + 0.4, 1.8, 0.3,
                     font_size=9, color=WHITE)

    add_text_box(s13, "自動打診フロー", 4.7, 1.4, 4.0, 0.5,
                 font_size=14, bold=True, color=PURPLE)
    flow_items = [
        "人員不足を検知",
        "グループ企業へ自動打診（LINE/メール）",
        "回答期限内に未回答 → 次順位へ自動エスカレーション",
        "承諾確認 → 配置に自動反映",
    ]
    for j, item in enumerate(flow_items):
        add_text_box(s13, f"{'↓' if j > 0 else '①'}  {item}",
                     4.7, 2.0 + j * 0.9, 4.0, 0.75,
                     font_size=10, color=DARK)

    add_text_box(s13, "融通実績管理", 9.2, 1.4, 3.8, 0.5,
                 font_size=14, bold=True, color=PURPLE)
    result_items = [
        "企業別融通回数・人数",
        "相互融通バランス可視化",
        "融通条件（単価・除外現場）",
        "月次レポート自動生成",
    ]
    for j, item in enumerate(result_items):
        add_text_box(s13, f"✓  {item}",
                     9.2, 2.0 + j * 0.9, 3.8, 0.75,
                     font_size=10, color=DARK)

    prs.save(path)
    print(f"PowerPoint更新完了: {path}")


if __name__ == "__main__":
    print("GuardSync プロジェクト計画書 v1.1 更新中...")
    print()
    update_word()
    update_excel()
    update_pptx()
    print()
    print("=" * 60)
    print("✓ 全ファイルの更新が完了しました（v1.1）")
    print(f"  保存先: {OUTPUT_DIR}")
