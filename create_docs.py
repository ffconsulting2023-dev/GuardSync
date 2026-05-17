#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
GuardSync プロジェクト計画書生成スクリプト
Word / Excel / PowerPoint の3形式で出力
"""

import os
from datetime import datetime, date

OUTPUT_DIR = "/Users/t.take/AI management/project/output/GuardSync"

# ============================================================
# WORD 生成
# ============================================================
def create_word():
    from docx import Document
    from docx.shared import Pt, RGBColor, Cm
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_TABLE_ALIGNMENT
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    doc = Document()

    # ページ余白設定
    for section in doc.sections:
        section.top_margin = Cm(2)
        section.bottom_margin = Cm(2)
        section.left_margin = Cm(2.5)
        section.right_margin = Cm(2.5)

    def add_heading(text, level=1, color=None):
        h = doc.add_heading(text, level=level)
        if color:
            for run in h.runs:
                run.font.color.rgb = RGBColor(*color)
        return h

    def add_table(headers, rows, col_widths=None):
        table = doc.add_table(rows=1, cols=len(headers))
        table.style = 'Table Grid'
        # ヘッダー行
        hdr = table.rows[0].cells
        for i, h in enumerate(headers):
            hdr[i].text = h
            hdr[i].paragraphs[0].runs[0].font.bold = True
            hdr[i].paragraphs[0].runs[0].font.size = Pt(9)
            hdr[i].paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER
        # データ行
        for row_data in rows:
            row = table.add_row().cells
            for i, val in enumerate(row_data):
                row[i].text = str(val)
                row[i].paragraphs[0].runs[0].font.size = Pt(9)
        return table

    # ====== 表紙 ======
    doc.add_paragraph()
    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run("GuardSync")
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1A, 0x56, 0xDB)

    subtitle = doc.add_paragraph()
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run2 = subtitle.add_run("警備会社向け統合管理システム\nプロジェクト計画書")
    run2.font.size = Pt(18)

    doc.add_paragraph()
    meta = doc.add_paragraph()
    meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
    meta.add_run(f"バージョン: 1.0　|　作成日: 2026年5月18日　|　ステータス: 承認待ち")

    doc.add_page_break()

    # ====== 0. エグゼクティブサマリー ======
    add_heading("0. エグゼクティブサマリー", 1, (0x1A, 0x56, 0xDB))
    doc.add_paragraph(
        "GuardSyncは、警備会社向け統合管理システムの新規開発プロジェクトです。"
        "競合サービス（komainu.cloud）の56機能を分析し、UI/UXの大幅な改善と"
        "発注元企業との自動連携・AI活用による差別化を実現します。\n\n"
        "MINプランを2026年6月20日にリリースし、MAXプランを2026年8月末に完成させます。"
        "開発はClaude Code（AI）を最大活用し、1名体制で実施します。"
    )

    # ====== 1. プロジェクト概要 ======
    add_heading("1. プロジェクト概要", 1, (0x1A, 0x56, 0xDB))
    add_heading("1.1 基本情報", 2)
    add_table(
        ["項目", "内容"],
        [
            ["プロジェクト名", "GuardSync"],
            ["サービス種別", "警備会社向け統合管理SaaS"],
            ["競合サービス", "komainu.cloud（56機能）"],
            ["開発体制", "1名 + Claude Code（AI自動実装）"],
            ["MINリリース", "2026年6月20日"],
            ["MAXリリース", "2026年8月末"],
            ["技術スタック", "React(PWA) + Express + Prisma + PostgreSQL + Railway"],
            ["インフラ", "Railway（バックエンド）+ お名前.com RS（フロントエンド）"],
        ]
    )

    doc.add_paragraph()
    add_heading("1.2 差別化ポイント", 2)
    doc.add_paragraph("競合（komainu.cloud）にない以下の機能で市場優位を確立します：")
    items = [
        "高品質UI/UX：高齢の警備員でも直感的に操作できるデザイン",
        "発注元連携：建設・建築・造園会社との双方向情報連携",
        "自動受付エンジン：LINE / メール / FAX（OCR）からの依頼を自動取得・サジェスト",
        "前日自動確認：翌日の依頼内容を前日午前中に自動送信・確認",
        "電子契約：独自実装の電子署名・タイムスタンプ・決裁フロー",
        "日払い管理：手数料管理・月末給与との自動差引",
    ]
    for item in items:
        p = doc.add_paragraph(style='List Bullet')
        p.add_run(item).font.size = Pt(10)

    # ====== 2. 要件定義（MoSCoW） ======
    doc.add_page_break()
    add_heading("2. 要件定義（MoSCoW分析）", 1, (0x1A, 0x56, 0xDB))

    add_heading("2.1 Must（6月20日必須）", 2)
    must_features = [
        ["管制機能", "状況監視・自動監視・1号/2号業務配置・月間シフト・ホワイトボード・管制表出力・業務レポート等（全15機能）"],
        ["隊員アプリPWA", "出動/上下番ワンボタン・月間スケジュール・経路案内・警備報告書（独自設計）"],
        ["マスターデータ", "隊員・現場/契約・協力会社・車両の各CRUD管理"],
        ["権限設定", "役割別アクセス制御（管制員・給与担当・請求担当等）"],
        ["マルチテナント", "複数警備会社が独立して利用できる設計"],
        ["請求機能", "請求作成・インボイス対応・総括請求・売掛管理・メール送信"],
        ["LINE Official連携", "お知らせ配信・アラート通知・既読管理"],
        ["自動受付エンジン", "LINE/メール/FAX→OCR→件名自動取得・作成サジェスト"],
        ["前日自動確認", "翌日配置を前日午前中に自動送信（差別化機能）"],
        ["実績管理", "打刻・時間記録・給与反映"],
        ["日払い管理", "申請/承認フロー・手数料管理・月末給与自動差引"],
        ["契約条件設定", "単価・勤務条件・請求への自動反映"],
        ["CSV出力", "勤怠・隊員・契約データのCSVダウンロード"],
    ]
    add_table(["カテゴリ", "機能内容"], must_features)

    doc.add_paragraph()
    add_heading("2.2 Should（v2：8月末）", 2)
    should_features = [
        ["給与集計", "日給月給・変形労働制・残業計算・給与明細メール送信"],
        ["労務管理", "就労時間アラート・有給管理・有給自動付与"],
        ["電子契約", "紙契約書OCR取込・社内決裁フロー・メール署名依頼・タイムスタンプ"],
        ["帳票強化", "帳票Excel出力・会計ソフトCSV（マネーフォワード/弥生/freee）"],
        ["隊員アプリ強化", "休暇申請・出勤希望・トイレ/コンビニ/駐車場検索"],
        ["発注元連携強化", "建設・造園会社との詳細情報連携"],
    ]
    add_table(["カテゴリ", "機能内容"], should_features)

    doc.add_paragraph()
    add_heading("2.3 Could（v3以降）", 2)
    could_features = [
        ["AI自動配置", "過去の配置傾向から自動シフト作成（競合も未実装）"],
        ["会計ソフトAPI連携", "マネーフォワード・freee・弥生への直接API連携"],
        ["IVR（自動音声電話）", "ガラケー対応の上下番自動音声報告"],
        ["日払い外部連携", "即給By GMO・JOBPAY連携"],
        ["年末調整", "年末調整書類の電子管理・計算"],
        ["認定タイムスタンプ", "セイコーソリューションズ等の認定TSA"],
    ]
    add_table(["カテゴリ", "機能内容"], could_features)

    add_heading("2.4 Won't（今回実施しない）", 2)
    wont_features = [
        ["IP電話連携", "専用機器・設定が複雑すぎる"],
        ["トイレ/コンビニ/駐車場検索", "地図API別途必要・優先度低（v2で検討）"],
        ["総括請求インボイス完全対応", "会計士確認が必要・複雑（v2で対応）"],
        ["有給給与計算", "社労士確認が必要（v2で対応）"],
    ]
    add_table(["機能", "理由"], wont_features)

    # ====== 3. スケジュール ======
    doc.add_page_break()
    add_heading("3. スケジュール（34日計画）", 1, (0x1A, 0x56, 0xDB))
    schedule = [
        ["Week 1", "5/17〜5/23", "基盤構築", "DB設計・認証・マルチテナント・マスターデータCRUD・権限設定・UIデザインシステム"],
        ["Week 2", "5/24〜5/30", "管制機能（前半）", "状況監視・自動監視・1号/2号業務配置・ホワイトボード"],
        ["Week 3", "5/31〜6/06", "管制機能（後半）＋隊員PWA", "月間シフト・管制表出力・業務レポート・隊員PWA（出動/スケジュール/経路/報告書）"],
        ["Week 4", "6/07〜6/13", "LINE連携・自動受付・請求", "LINE Official・前日確認・メール/FAX自動受付・OCR・請求作成・売掛管理"],
        ["Week 5", "6/14〜6/20", "日払い・実績・テスト・デプロイ", "日払い管理・実績管理・CSV出力・結合テスト・UI調整・本番デプロイ"],
    ]
    add_table(["週", "期間", "フェーズ", "主なタスク"], schedule)

    doc.add_paragraph()
    add_heading("3.1 マイルストーン", 2)
    milestones = [
        ["M1", "2026/05/23", "基盤完成", "認証・マルチテナント・全マスターデータCRUD動作確認"],
        ["M2", "2026/05/30", "管制機能前半完成", "状況監視・配置機能の動作確認"],
        ["M3", "2026/06/06", "管制機能完成＋PWA完成", "全管制機能・隊員PWAの動作確認"],
        ["M4", "2026/06/13", "連携・請求完成", "LINE連携・自動受付・請求機能の動作確認"],
        ["M5", "2026/06/20", "MINリリース", "全Must機能テスト完了・本番デプロイ・サービス開始"],
        ["M6", "2026/08/31", "MAXリリース", "全Should機能テスト完了・v2リリース"],
    ]
    add_table(["#", "目標日", "マイルストーン", "完了条件"], milestones)

    # ====== 4. 技術構成 ======
    doc.add_page_break()
    add_heading("4. 技術構成", 1, (0x1A, 0x56, 0xDB))
    tech = [
        ["フロントエンド（管理画面）", "React + TypeScript + Vite + TailwindCSS"],
        ["フロントエンド（隊員アプリ）", "React PWA + Service Worker（ホーム画面追加対応）"],
        ["バックエンド", "Express.js + TypeScript"],
        ["ORM / DB", "Prisma + PostgreSQL"],
        ["インフラ", "Railway（バックエンド・DB）+ お名前.com RS（フロントエンド）"],
        ["メッセージング", "LINE Official Account + Messaging API"],
        ["OCR", "Google Vision API（FAX・紙契約書の文字認識）"],
        ["PDF生成", "PDFKit or Puppeteer（請求書・管制表・警備報告書）"],
        ["メール送信", "SendGrid or Nodemailer + SMTP"],
        ["電子署名", "独自実装（RFC 3161タイムスタンプ + 合意証明型）"],
        ["認証", "JWT + bcrypt"],
        ["スケジューラー", "node-cron（前日自動確認・アラート）"],
    ]
    add_table(["技術要素", "採用技術"], tech)

    # ====== 5. リスク管理 ======
    doc.add_page_break()
    add_heading("5. リスク管理", 1, (0x1A, 0x56, 0xDB))
    risks = [
        ["R1", "34日でのスコープ超過", "高", "大", "🔴高", "毎週末にMust機能の進捗確認・Should機能をv2へ移動"],
        ["R2", "LINE Official審査遅延", "中", "大", "🔴高", "今週中に申請・審査中はモック実装で並行開発"],
        ["R3", "警備報告書の特許侵害", "低", "大", "🟡中", "独自設計（QRコード承認方式）で完全回避"],
        ["R4", "給与計算ロジックの不正確", "中", "大", "🔴高", "基本形態のみMINに含め・社労士確認をv2実装前に実施"],
        ["R5", "FAX→メール変換サービスの設定遅延", "中", "中", "🟡中", "Week1中に契約完了・設定はWeek4前に完了させる"],
        ["R6", "Google Vision API費用超過", "低", "小", "🟢低", "月間コスト上限を設定・キャッシュ実装で重複OCRを防止"],
    ]
    add_table(["#", "リスク", "発生確率", "影響度", "優先度", "対応策"], risks)

    # ====== 6. セキュリティ ======
    doc.add_page_break()
    add_heading("6. セキュリティ方針", 1, (0x1A, 0x56, 0xDB))
    security = [
        ["認証・認可", "JWT認証・ロールベースアクセス制御・セッション管理"],
        ["データ保護", "テナント間完全分離・SQLインジェクション対策（Prisma ORM）"],
        ["通信", "HTTPS強制・HSTS設定"],
        ["個人情報", "隊員の個人情報・位置情報は暗号化保存"],
        ["電子契約", "タイムスタンプ・ハッシュ値による改ざん検知・監査ログ"],
        ["パスワード", "bcryptハッシュ化・ブルートフォース対策"],
        ["外部API", "APIキーを環境変数管理・ソースコードに含めない"],
    ]
    add_table(["対策項目", "内容"], security)

    # ====== 7. 事前タスク ======
    add_heading("7. 開発開始前の事前タスク", 1, (0x1A, 0x56, 0xDB))
    pretasks = [
        ["LINE Official Account 開設", "2026/05/18", "審査に数日かかる場合あり・今日中に申請"],
        ["Messaging API チャネル作成", "2026/05/18", "LINE Developer Consoleで作成"],
        ["FAX→メール変換サービス契約", "2026/05/20", "eFax・メールdeFA X等"],
        ["Claude Code Maxプラン加入", "2026/05/18", "開発効率最大化のため"],
        ["GitHubリポジトリ作成", "2026/05/18", "guardsync リポジトリを新規作成"],
        ["Railway新規プロジェクト作成", "2026/05/18", "本番環境の準備"],
    ]
    add_table(["タスク", "期限", "備考"], pretasks)

    # ====== フッター ======
    doc.add_paragraph()
    footer_p = doc.add_paragraph()
    footer_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    footer_run = footer_p.add_run("GuardSync プロジェクト計画書 v1.0 | 2026年5月18日")
    footer_run.font.size = Pt(8)
    footer_run.font.color.rgb = RGBColor(0x88, 0x88, 0x88)

    path = os.path.join(OUTPUT_DIR, "20260518_GuardSync_プロジェクト計画書.docx")
    doc.save(path)
    print(f"Word出力完了: {path}")
    return path


# ============================================================
# EXCEL 生成
# ============================================================
def create_excel():
    from openpyxl import Workbook
    from openpyxl.styles import (
        Font, PatternFill, Alignment, Border, Side, GradientFill
    )
    from openpyxl.utils import get_column_letter

    wb = Workbook()

    NAVY = "1A56DB"
    LIGHT_BLUE = "DBEAFE"
    GREEN = "16A34A"
    LIGHT_GREEN = "DCFCE7"
    YELLOW = "CA8A04"
    LIGHT_YELLOW = "FEF9C3"
    RED = "DC2626"
    LIGHT_RED = "FEE2E2"
    GRAY = "6B7280"
    LIGHT_GRAY = "F3F4F6"
    WHITE = "FFFFFF"

    header_font = Font(bold=True, color=WHITE, size=10)
    header_fill = PatternFill("solid", fgColor=NAVY)
    header_align = Alignment(horizontal="center", vertical="center", wrap_text=True)
    body_font = Font(size=9)
    body_align = Alignment(vertical="center", wrap_text=True)
    center_align = Alignment(horizontal="center", vertical="center", wrap_text=True)
    thin_border = Border(
        left=Side(style="thin", color="CCCCCC"),
        right=Side(style="thin", color="CCCCCC"),
        top=Side(style="thin", color="CCCCCC"),
        bottom=Side(style="thin", color="CCCCCC"),
    )

    def style_header(ws, row, cols):
        for col in range(1, cols + 1):
            c = ws.cell(row=row, column=col)
            c.font = header_font
            c.fill = header_fill
            c.alignment = header_align
            c.border = thin_border

    def style_row(ws, row, cols, fill_color=None, center_cols=None):
        center_cols = center_cols or []
        for col in range(1, cols + 1):
            c = ws.cell(row=row, column=col)
            c.font = body_font
            if fill_color:
                c.fill = PatternFill("solid", fgColor=fill_color)
            c.alignment = center_align if col in center_cols else body_align
            c.border = thin_border

    # ====== Sheet1: 概要 ======
    ws1 = wb.active
    ws1.title = "概要"

    ws1.merge_cells("A1:F1")
    ws1["A1"] = "GuardSync プロジェクト計画書"
    ws1["A1"].font = Font(bold=True, size=18, color=NAVY)
    ws1["A1"].alignment = Alignment(horizontal="center", vertical="center")
    ws1["A1"].fill = PatternFill("solid", fgColor=LIGHT_BLUE)
    ws1.row_dimensions[1].height = 40

    ws1.merge_cells("A2:F2")
    ws1["A2"] = "警備会社向け統合管理システム | バージョン1.0 | 2026年5月18日"
    ws1["A2"].font = Font(size=10, color=GRAY)
    ws1["A2"].alignment = Alignment(horizontal="center")

    ws1.row_dimensions[3].height = 10

    # 基本情報
    ws1["A4"] = "基本情報"
    ws1["A4"].font = Font(bold=True, size=11, color=NAVY)
    info = [
        ("プロジェクト名", "GuardSync"),
        ("サービス種別", "警備会社向け統合管理SaaS"),
        ("競合", "komainu.cloud（56機能）"),
        ("開発体制", "1名 + Claude Code（AI）"),
        ("MINリリース", "2026年6月20日"),
        ("MAXリリース", "2026年8月末"),
        ("技術スタック", "React(PWA) + Express + Prisma + PostgreSQL + Railway"),
    ]
    for i, (k, v) in enumerate(info, start=5):
        ws1.cell(row=i, column=1, value=k).font = Font(bold=True, size=9)
        ws1.cell(row=i, column=1).fill = PatternFill("solid", fgColor=LIGHT_BLUE)
        ws1.cell(row=i, column=1).border = thin_border
        ws1.cell(row=i, column=1).alignment = center_align
        ws1.merge_cells(f"B{i}:F{i}")
        ws1.cell(row=i, column=2, value=v).font = Font(size=9)
        ws1.cell(row=i, column=2).border = thin_border

    ws1.column_dimensions["A"].width = 20
    ws1.column_dimensions["B"].width = 60

    # ====== Sheet2: WBS・タスク ======
    ws2 = wb.create_sheet("WBS・タスク")

    headers = ["#", "週", "期間", "フェーズ", "タスク", "工数(日)", "依存", "担当", "完了条件", "ステータス"]
    for col, h in enumerate(headers, 1):
        ws2.cell(row=1, column=col, value=h)
    style_header(ws2, 1, len(headers))
    ws2.row_dimensions[1].height = 25

    tasks = [
        # Week1
        ("1.1", "Week1", "5/17-5/23", "基盤構築", "新規リポジトリ・Railway環境構築", 0.5, "-", "開発者", "デプロイ疎通確認", "未着手"),
        ("1.2", "Week1", "5/17-5/23", "基盤構築", "Prisma DBスキーマ設計（全テーブル）", 1.5, "-", "開発者", "マイグレーション完了", "未着手"),
        ("1.3", "Week1", "5/17-5/23", "基盤構築", "認証実装（JWT・セッション管理）", 0.5, "1.2", "開発者", "ログイン/ログアウト動作", "未着手"),
        ("1.4", "Week1", "5/17-5/23", "基盤構築", "マルチテナント実装（会社ID分離）", 0.5, "1.3", "開発者", "テナント間データ分離確認", "未着手"),
        ("1.5", "Week1", "5/17-5/23", "基盤構築", "権限設定（ロール・アクセス制御）", 0.5, "1.4", "開発者", "役割別アクセス制御動作", "未着手"),
        ("1.6", "Week1", "5/17-5/23", "基盤構築", "隊員情報CRUD", 0.5, "1.2", "開発者", "一覧・登録・編集・削除", "未着手"),
        ("1.7", "Week1", "5/17-5/23", "基盤構築", "現場/契約情報CRUD（緯度経度含む）", 0.5, "1.2", "開発者", "一覧・登録・編集・削除", "未着手"),
        ("1.8", "Week1", "5/17-5/23", "基盤構築", "協力会社・車両情報CRUD", 0.5, "1.2", "開発者", "一覧・登録・編集・削除", "未着手"),
        ("1.9", "Week1", "5/17-5/23", "基盤構築", "UIデザインシステム（共通コンポーネント）", 0.5, "-", "開発者", "共通ヘッダー・ナビ・テーブル完成", "未着手"),
        # Week2
        ("2.1", "Week2", "5/24-5/30", "管制（前半）", "状況監視画面（リアルタイム・GPS）", 1.5, "1.2", "開発者", "ステータス自動更新・位置情報表示", "未着手"),
        ("2.2", "Week2", "5/24-5/30", "管制（前半）", "自動監視（遅刻アラート・未対応フラグ）", 1.0, "2.1", "開発者", "指定時間超過でアラート発生", "未着手"),
        ("2.3", "Week2", "5/24-5/30", "管制（前半）", "1号業務配置（月次シフト作成）", 1.5, "1.7", "開発者", "カレンダー形式でシフト作成・印刷", "未着手"),
        ("2.4", "Week2", "5/24-5/30", "管制（前半）", "2号業務配置（日・週次・前日コピー）", 1.5, "1.7", "開発者", "現場カード形式・前日コピー動作", "未着手"),
        ("2.5", "Week2", "5/24-5/30", "管制（前半）", "ホワイトボード（全現場俯瞰）", 1.0, "2.1", "開発者", "リアルタイムステータスバッジ表示", "未着手"),
        # Week3
        ("3.1", "Week3", "5/31-6/6", "管制（後半）", "月間シフト（隊員主語・稼働率集計）", 1.0, "2.3", "開発者", "隊員別カレンダー・稼働率表示", "未着手"),
        ("3.2", "Week3", "5/31-6/6", "管制（後半）", "月間受注状況（顧客×現場カレンダー）", 0.5, "1.7", "開発者", "日別人工数・受注額集計", "未着手"),
        ("3.3", "Week3", "5/31-6/6", "管制（後半）", "経営データダッシュボード（粗利・目標）", 1.0, "3.2", "開発者", "顧客別粗利・達成率チャート", "未着手"),
        ("3.4", "Week3", "5/31-6/6", "管制（後半）", "管制表PDF出力（テンプレート）", 0.5, "2.3", "開発者", "PDF生成・印刷動作", "未着手"),
        ("3.5", "Week3", "5/31-6/6", "管制（後半）", "警備報告書（独自設計・電子サイン）", 0.5, "1.6", "開発者", "報告書生成・サイン取得・PDF出力", "未着手"),
        ("3.6", "Week3", "5/31-6/6", "管制（後半）", "業務レポート（未完了タスク一覧）", 0.5, "2.1", "開発者", "未対応項目の自動ピックアップ", "未着手"),
        ("3.7", "Week3", "5/31-6/6", "隊員PWA", "隊員PWA基盤（manifest・SW設定）", 0.5, "1.3", "開発者", "ホーム画面追加・オフライン対応", "未着手"),
        ("3.8", "Week3", "5/31-6/6", "隊員PWA", "出動・上下番ワンボタン（GPS打刻）", 0.5, "3.7", "開発者", "ワンタップで状態遷移・管制側反映", "未着手"),
        ("3.9", "Week3", "5/31-6/6", "隊員PWA", "月間スケジュール（PWA・色分け）", 0.5, "3.7", "開発者", "カレンダー表示・現場詳細タップ", "未着手"),
        ("3.10", "Week3", "5/31-6/6", "隊員PWA", "経路案内（マップdeeplink）", 0.5, "3.7", "開発者", "現場住所からナビアプリ起動", "未着手"),
        # Week4
        ("4.1", "Week4", "6/7-6/13", "LINE・自動受付", "LINE Official設定・Messaging API連携", 0.5, "-", "開発者", "Webhook疎通・メッセージ送受信", "未着手"),
        ("4.2", "Week4", "6/7-6/13", "LINE・自動受付", "お知らせ配信（既読管理・LINE送信）", 0.5, "4.1", "開発者", "配信・既読確認・返信記録", "未着手"),
        ("4.3", "Week4", "6/7-6/13", "LINE・自動受付", "アラートLINE通知（遅刻・未報告）", 0.5, "4.1", "開発者", "条件発火→LINE自動送信", "未着手"),
        ("4.4", "Week4", "6/7-6/13", "LINE・自動受付", "前日午前中 自動確認送信（Cron）", 0.5, "4.1", "開発者", "毎朝10時に翌日配置を自動送信", "未着手"),
        ("4.5", "Week4", "6/7-6/13", "LINE・自動受付", "LINE Webhook自動受付→サジェスト", 1.0, "4.1", "開発者", "受信→内容解析→件名サジェスト表示", "未着手"),
        ("4.6", "Week4", "6/7-6/13", "LINE・自動受付", "メール自動受付→件名取得・サジェスト", 1.0, "-", "開発者", "メール受信→解析→サジェスト表示", "未着手"),
        ("4.7", "Week4", "6/7-6/13", "LINE・自動受付", "FAX→PDF受信→OCR→サジェスト", 1.0, "-", "開発者", "PDF受信→OCR→テキスト化→サジェスト", "未着手"),
        ("4.8", "Week4", "6/7-6/13", "請求", "請求作成（印鑑・インボイス・集計）", 1.0, "1.7", "開発者", "ワンクリック集計・PDF生成・印鑑合成", "未着手"),
        ("4.9", "Week4", "6/7-6/13", "請求", "総括請求・グルーピング", 0.5, "4.8", "開発者", "複数現場を1枚にまとめ出力", "未着手"),
        ("4.10", "Week4", "6/7-6/13", "請求", "売掛管理（入金・残高管理）", 0.5, "4.8", "開発者", "顧客別売掛一覧・残高確認", "未着手"),
        ("4.11", "Week4", "6/7-6/13", "請求", "請求書メール自動送信（送信履歴）", 0.5, "4.8", "開発者", "PDF添付メール送信・履歴記録", "未着手"),
        # Week5
        ("5.1", "Week5", "6/14-6/20", "仕上げ", "実績管理（打刻・時間記録・給与反映）", 0.5, "2.1", "開発者", "予定/打刻/報告書の3データ記録", "未着手"),
        ("5.2", "Week5", "6/14-6/20", "仕上げ", "日払い申請・承認フロー（PWA/LINE）", 0.5, "4.1", "開発者", "申請→承認→記録の動作確認", "未着手"),
        ("5.3", "Week5", "6/14-6/20", "仕上げ", "手数料管理（固定/率・負担者設定）", 0.5, "5.2", "開発者", "手数料パターン設定・計算確認", "未着手"),
        ("5.4", "Week5", "6/14-6/20", "仕上げ", "月末給与自動差引（日払い+手数料）", 0.5, "5.3", "開発者", "差引計算・給与明細への明記", "未着手"),
        ("5.5", "Week5", "6/14-6/20", "仕上げ", "契約条件設定（単価・勤務条件）", 0.5, "1.7", "開発者", "条件設定→請求自動反映", "未着手"),
        ("5.6", "Week5", "6/14-6/20", "仕上げ", "各種CSV出力（勤怠・隊員・契約）", 0.5, "-", "開発者", "各データをCSVダウンロード", "未着手"),
        ("5.7", "Week5", "6/14-6/20", "仕上げ", "結合テスト・バグ修正", 1.0, "-", "開発者", "主要フロー全て動作確認", "未着手"),
        ("5.8", "Week5", "6/14-6/20", "仕上げ", "UI/UXブラッシュアップ", 0.5, "-", "開発者", "デザイン品質確認・レスポンシブ対応", "未着手"),
        ("5.9", "Week5", "6/14-6/20", "仕上げ", "本番デプロイ・動作確認", 0.5, "5.7", "開発者", "Railway本番環境で全機能動作", "未着手"),
    ]

    week_colors = {
        "Week1": "DBEAFE", "Week2": "DCFCE7",
        "Week3": "FEF9C3", "Week4": "FFE4E6", "Week5": "F3E8FF"
    }
    status_colors = {"未着手": None, "進行中": LIGHT_YELLOW, "完了": LIGHT_GREEN, "遅延": LIGHT_RED}

    for i, task in enumerate(tasks, start=2):
        for col, val in enumerate(task, start=1):
            c = ws2.cell(row=i, column=col, value=val)
            c.font = body_font
            c.border = thin_border
            c.alignment = center_align if col in [1, 2, 3, 6, 7, 8, 10] else body_align
        fill_color = week_colors.get(task[1], WHITE)
        for col in range(1, len(headers) + 1):
            ws2.cell(row=i, column=col).fill = PatternFill("solid", fgColor=fill_color)
        ws2.row_dimensions[i].height = 30

    col_widths = [6, 8, 12, 15, 35, 8, 8, 10, 30, 10]
    for col, width in enumerate(col_widths, 1):
        ws2.column_dimensions[get_column_letter(col)].width = width
    ws2.freeze_panes = "A2"

    # ====== Sheet3: スケジュール ======
    ws3 = wb.create_sheet("スケジュール・マイルストーン")

    ms_headers = ["#", "マイルストーン", "目標日", "完了条件", "ステータス"]
    for col, h in enumerate(ms_headers, 1):
        ws3.cell(row=1, column=col, value=h)
    style_header(ws3, 1, len(ms_headers))

    milestones = [
        ("M1", "基盤完成", "2026/05/23", "認証・マルチテナント・全マスターデータCRUD動作確認", "未着手"),
        ("M2", "管制機能前半完成", "2026/05/30", "状況監視・1号/2号配置・ホワイトボードの動作確認", "未着手"),
        ("M3", "管制機能＋隊員PWA完成", "2026/06/06", "全管制機能・隊員PWA基本4機能の動作確認", "未着手"),
        ("M4", "連携・請求完成", "2026/06/13", "LINE連携・自動受付・請求機能の動作確認", "未着手"),
        ("M5", "MINリリース", "2026/06/20", "全Must機能テスト完了・本番デプロイ・サービス開始", "未着手"),
        ("M6", "MAXリリース", "2026/08/31", "全Should機能テスト完了・v2本番リリース", "未着手"),
    ]
    for i, ms in enumerate(milestones, start=2):
        for col, val in enumerate(ms, start=1):
            c = ws3.cell(row=i, column=col, value=val)
            c.font = body_font
            c.border = thin_border
            c.alignment = center_align if col in [1, 3, 5] else body_align
        color = LIGHT_GREEN if ms[0] == "M5" or ms[0] == "M6" else LIGHT_BLUE
        for col in range(1, len(ms_headers) + 1):
            ws3.cell(row=i, column=col).fill = PatternFill("solid", fgColor=color)
        ws3.row_dimensions[i].height = 35

    ws3.column_dimensions["A"].width = 6
    ws3.column_dimensions["B"].width = 25
    ws3.column_dimensions["C"].width = 14
    ws3.column_dimensions["D"].width = 50
    ws3.column_dimensions["E"].width = 10

    # ====== Sheet4: リスク管理 ======
    ws4 = wb.create_sheet("リスク管理")

    risk_headers = ["#", "リスク", "発生確率", "影響度", "優先度", "予防策", "発生時対応"]
    for col, h in enumerate(risk_headers, 1):
        ws4.cell(row=1, column=col, value=h)
    style_header(ws4, 1, len(risk_headers))

    risks = [
        ("R1", "34日でのスコープ超過", "高", "大", "🔴高", "毎週末に進捗確認・Should機能をv2へ移動", "Must機能を優先・Should機能をスコープ外へ"),
        ("R2", "LINE Official審査遅延", "中", "大", "🔴高", "今週中に申請・モック実装で並行開発", "モック実装で他機能開発を継続"),
        ("R3", "警備報告書の特許侵害", "低", "大", "🟡中", "独自設計（QRコード承認方式）で完全回避", "設計を変更・弁護士確認"),
        ("R4", "給与計算ロジックの不正確", "中", "大", "🔴高", "基本形態のみMINに含め社労士確認", "v2実装前に社労士レビュー必須"),
        ("R5", "FAX→メール変換設定遅延", "中", "中", "🟡中", "Week1中に契約・設定完了", "メール/LINE受付のみで代替"),
        ("R6", "Google Vision API費用超過", "低", "小", "🟢低", "月間コスト上限設定・キャッシュ実装", "OCR回数制限・手動入力を代替手段に"),
    ]
    risk_colors = {"🔴高": LIGHT_RED, "🟡中": LIGHT_YELLOW, "🟢低": LIGHT_GREEN}
    for i, risk in enumerate(risks, start=2):
        for col, val in enumerate(risk, start=1):
            c = ws4.cell(row=i, column=col, value=val)
            c.font = body_font
            c.border = thin_border
            c.alignment = center_align if col in [1, 3, 4, 5] else body_align
        fill = risk_colors.get(risk[4], WHITE)
        for col in range(1, len(risk_headers) + 1):
            ws4.cell(row=i, column=col).fill = PatternFill("solid", fgColor=fill)
        ws4.row_dimensions[i].height = 40

    ws4.column_dimensions["A"].width = 6
    ws4.column_dimensions["B"].width = 25
    ws4.column_dimensions["C"].width = 10
    ws4.column_dimensions["D"].width = 8
    ws4.column_dimensions["E"].width = 8
    ws4.column_dimensions["F"].width = 35
    ws4.column_dimensions["G"].width = 35

    # ====== Sheet5: ToDo ======
    ws5 = wb.create_sheet("ToDo・事前タスク")

    todo_headers = ["優先度", "期限", "タスク", "担当", "備考", "完了"]
    for col, h in enumerate(todo_headers, 1):
        ws5.cell(row=1, column=col, value=h)
    style_header(ws5, 1, len(todo_headers))

    todos = [
        ("🔴緊急", "2026/05/18", "LINE Official Account 開設申請", "社長", "審査に数日かかる場合あり", "☐"),
        ("🔴緊急", "2026/05/18", "Claude Code Maxプラン加入", "社長", "claude.ai/settings から加入", "☐"),
        ("🔴緊急", "2026/05/18", "GitHubリポジトリ作成（guardsync）", "社長", "新規リポジトリ", "☐"),
        ("🔴緊急", "2026/05/18", "Railway新規プロジェクト作成", "社長", "本番環境準備", "☐"),
        ("🟡重要", "2026/05/20", "FAX→メール変換サービス契約", "社長", "eFax等・月額費用確認", "☐"),
        ("🟡重要", "2026/05/20", "LINE Messaging APIチャネル作成", "社長", "LINE Developers Consoleで設定", "☐"),
        ("🟡重要", "2026/05/20", "Google Vision API有効化", "社長", "Google Cloud Console・課金設定", "☐"),
        ("🟢通常", "2026/05/23", "SendGrid or SMTP設定", "開発者", "メール送信用", "☐"),
        ("🟢通常", "2026/05/23", "Railway DB（PostgreSQL）作成", "開発者", "本番DB環境", "☐"),
    ]
    priority_colors = {"🔴緊急": LIGHT_RED, "🟡重要": LIGHT_YELLOW, "🟢通常": LIGHT_GREEN}
    for i, todo in enumerate(todos, start=2):
        for col, val in enumerate(todo, start=1):
            c = ws5.cell(row=i, column=col, value=val)
            c.font = body_font
            c.border = thin_border
            c.alignment = center_align if col in [1, 2, 4, 6] else body_align
        fill = priority_colors.get(todo[0], WHITE)
        for col in range(1, len(todo_headers) + 1):
            ws5.cell(row=i, column=col).fill = PatternFill("solid", fgColor=fill)
        ws5.row_dimensions[i].height = 25

    ws5.column_dimensions["A"].width = 10
    ws5.column_dimensions["B"].width = 14
    ws5.column_dimensions["C"].width = 40
    ws5.column_dimensions["D"].width = 10
    ws5.column_dimensions["E"].width = 30
    ws5.column_dimensions["F"].width = 8

    # 全シートのフィルター設定
    for ws in [ws2, ws3, ws4, ws5]:
        ws.auto_filter.ref = ws.dimensions

    path = os.path.join(OUTPUT_DIR, "20260518_GuardSync_プロジェクト計画書.xlsx")
    wb.save(path)
    print(f"Excel出力完了: {path}")
    return path


# ============================================================
# POWERPOINT 生成
# ============================================================
def create_pptx():
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    NAVY = RGBColor(0x1A, 0x56, 0xDB)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT_BLUE = RGBColor(0xDB, 0xEA, 0xFE)
    DARK = RGBColor(0x1F, 0x29, 0x37)
    GRAY = RGBColor(0x6B, 0x72, 0x80)
    GREEN = RGBColor(0x16, 0xA3, 0x4A)
    RED = RGBColor(0xDC, 0x26, 0x26)
    YELLOW = RGBColor(0xCA, 0x8A, 0x04)

    blank_layout = prs.slide_layouts[6]

    def add_rect(slide, x, y, w, h, fill_color, line_color=None):
        shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        if line_color:
            shape.line.color.rgb = line_color
        else:
            shape.line.fill.background()
        return shape

    def add_text_box(slide, text, x, y, w, h, font_size=14, bold=False,
                     color=None, align=PP_ALIGN.LEFT, italic=False):
        txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        if color:
            run.font.color.rgb = color
        return txBox

    def slide_header(slide, title, subtitle=None):
        add_rect(slide, 0, 0, 13.33, 1.2, NAVY)
        add_text_box(slide, title, 0.4, 0.15, 12, 0.7,
                     font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
        if subtitle:
            add_text_box(slide, subtitle, 0.4, 0.75, 12, 0.4,
                         font_size=13, color=LIGHT_BLUE, align=PP_ALIGN.LEFT)
        add_rect(slide, 0, 7.3, 13.33, 0.2, NAVY)
        add_text_box(slide, "GuardSync プロジェクト計画書 v1.0 | 2026.05.18",
                     0.4, 7.3, 12, 0.2, font_size=8, color=WHITE, align=PP_ALIGN.CENTER)

    # ====== スライド1: 表紙 ======
    s1 = prs.slides.add_slide(blank_layout)
    add_rect(s1, 0, 0, 13.33, 7.5, NAVY)
    add_rect(s1, 0, 4.5, 13.33, 3.0, RGBColor(0x0F, 0x33, 0x8A))
    add_text_box(s1, "GuardSync", 1, 1.5, 11, 1.5,
                 font_size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(s1, "警備会社向け統合管理システム", 1, 2.9, 11, 0.8,
                 font_size=24, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
    add_text_box(s1, "プロジェクト計画書 v1.0", 1, 3.6, 11, 0.6,
                 font_size=16, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
    add_text_box(s1, "2026年5月18日　|　MINリリース: 6月20日　|　MAXリリース: 8月末",
                 1, 5.5, 11, 0.5, font_size=13, color=WHITE, align=PP_ALIGN.CENTER)

    # ====== スライド2: アジェンダ ======
    s2 = prs.slides.add_slide(blank_layout)
    slide_header(s2, "アジェンダ")
    agenda = [
        ("01", "プロジェクト概要", "背景・目的・チーム構成"),
        ("02", "競合分析", "komainu.cloud（56機能）との比較"),
        ("03", "差別化戦略", "4つの競合優位ポイント"),
        ("04", "要件定義", "MoSCoW分析・Must機能一覧"),
        ("05", "スケジュール", "34日計画・マイルストーン"),
        ("06", "技術構成", "アーキテクチャ・使用技術"),
        ("07", "リスク管理", "主要リスクと対応策"),
        ("08", "次のアクション", "今日からやること"),
    ]
    cols = [(0.4, 1.4), (6.8, 1.4)]
    for i, (num, title, sub) in enumerate(agenda):
        col_idx = i % 2
        row_idx = i // 2
        x = cols[col_idx][0]
        y = 1.3 + row_idx * 1.35
        add_rect(s2, x, y, 6.0, 1.1, LIGHT_BLUE)
        add_text_box(s2, num, x + 0.1, y + 0.1, 0.6, 0.8,
                     font_size=20, bold=True, color=NAVY)
        add_text_box(s2, title, x + 0.7, y + 0.05, 5.0, 0.5,
                     font_size=14, bold=True, color=DARK)
        add_text_box(s2, sub, x + 0.7, y + 0.55, 5.0, 0.4,
                     font_size=10, color=GRAY)

    # ====== スライド3: プロジェクト概要 ======
    s3 = prs.slides.add_slide(blank_layout)
    slide_header(s3, "01. プロジェクト概要", "GuardSyncとは何か")
    info_items = [
        ("プロジェクト名", "GuardSync"),
        ("ターゲット", "警備会社（全国）"),
        ("競合", "komainu.cloud（56機能）"),
        ("開発体制", "1名 + Claude Code（AI）毎日フル稼働"),
        ("MINリリース", "2026年6月20日（34日後）"),
        ("MAXリリース", "2026年8月末"),
        ("技術", "React PWA + Express + Prisma + Railway"),
    ]
    for i, (k, v) in enumerate(info_items):
        row = i % 4
        col_set = i // 4
        x = 0.4 + col_set * 6.4
        y = 1.4 + row * 1.35
        if col_set == 0 and i < 4:
            add_rect(s3, x, y, 6.0, 1.1, LIGHT_BLUE)
            add_text_box(s3, k, x + 0.2, y + 0.05, 5.6, 0.4,
                         font_size=10, color=GRAY)
            add_text_box(s3, v, x + 0.2, y + 0.45, 5.6, 0.5,
                         font_size=14, bold=True, color=DARK)
        elif col_set == 1 and i < 7:
            add_rect(s3, x, y, 6.4, 1.1, LIGHT_BLUE)
            add_text_box(s3, k, x + 0.2, y + 0.05, 6.0, 0.4,
                         font_size=10, color=GRAY)
            add_text_box(s3, v, x + 0.2, y + 0.45, 6.0, 0.5,
                         font_size=14, bold=True, color=DARK)

    # ====== スライド4: 競合分析 ======
    s4 = prs.slides.add_slide(blank_layout)
    slide_header(s4, "02. 競合分析", "komainu.cloud の機能構成")
    comp_data = [
        ("管制機能", "15機能", "状況監視・配置・シフト・報告書・管制表など"),
        ("隊員アプリ", "11機能", "出動連絡・スケジュール・チャット・休暇申請など"),
        ("請求機能", "4機能", "請求作成・インボイス・売掛管理・メール送信"),
        ("給与機能", "7機能", "実績管理・給与集計・年末調整・日払い連携など"),
        ("労務管理", "4機能", "有給自動付与・有給管理・就労時間アラートなど"),
        ("帳票/データ連携", "8機能", "CSV/Excel出力・勤怠データ・給与書類など"),
        ("便利機能", "7機能", "権限設定・チャット・アラート・IVRなど"),
    ]
    add_rect(s4, 0.4, 1.3, 12.5, 0.4, NAVY)
    add_text_box(s4, "競合機能数合計：56機能", 0.4, 1.3, 12.5, 0.4,
                 font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for i, (cat, count, detail) in enumerate(comp_data):
        row = i % 4
        col_s = i // 4
        x = 0.4 + col_s * 6.5
        y = 1.85 + row * 1.3
        w = 6.0 if col_s == 0 else 6.4
        add_rect(s4, x, y, w, 1.1, RGBColor(0xF0, 0xF4, 0xFF))
        add_text_box(s4, cat, x + 0.15, y + 0.05, w - 2.5, 0.4,
                     font_size=12, bold=True, color=DARK)
        add_text_box(s4, count, x + w - 1.2, y + 0.05, 1.0, 0.4,
                     font_size=14, bold=True, color=NAVY, align=PP_ALIGN.RIGHT)
        add_text_box(s4, detail, x + 0.15, y + 0.5, w - 0.3, 0.5,
                     font_size=9, color=GRAY)

    # ====== スライド5: 差別化戦略 ======
    s5 = prs.slides.add_slide(blank_layout)
    slide_header(s5, "03. 差別化戦略", "競合にない4つの優位ポイント")
    diff_items = [
        ("🎨", "UI/UX", "高品質なデザインと直感的な操作性\n高齢の警備員でもワンタップで操作完結"),
        ("🔗", "発注元連携", "建設・建築・造園会社との双方向情報連携\n受注から配置まで自動フロー化"),
        ("🤖", "自動受付エンジン", "LINE / メール / FAX（OCR）からの依頼を\n自動取得・件名生成サジェスト"),
        ("📅", "前日自動確認", "翌日の依頼内容を前日午前中に自動送信\nゼロアクションで抜け漏れ防止"),
    ]
    add_rect(s5, 0.3, 1.3, 12.7, 5.8, RGBColor(0xF8, 0xFA, 0xFF))
    for i, (icon, title, desc) in enumerate(diff_items):
        x = 0.5 + (i % 2) * 6.4
        y = 1.5 + (i // 2) * 2.8
        add_rect(s5, x, y, 6.0, 2.4, NAVY)
        add_text_box(s5, icon, x + 0.2, y + 0.2, 0.8, 0.8, font_size=28, color=WHITE)
        add_text_box(s5, title, x + 1.1, y + 0.2, 4.6, 0.6,
                     font_size=18, bold=True, color=WHITE)
        add_text_box(s5, desc, x + 0.2, y + 1.0, 5.5, 1.2,
                     font_size=11, color=LIGHT_BLUE)

    # ====== スライド6: MoSCoW（Must） ======
    s6 = prs.slides.add_slide(blank_layout)
    slide_header(s6, "04. 要件定義（Must機能）", "6月20日リリースに含む機能")
    must_items = [
        "管制機能 全15機能（状況監視・配置・シフト・管制表・業務レポート等）",
        "隊員アプリPWA（出動/上下番・月間スケジュール・経路案内・警備報告書）",
        "マスターデータ（隊員・現場/契約・協力会社・車両）",
        "権限設定・マルチテナント（複数警備会社対応）",
        "請求機能（請求作成・インボイス・総括・売掛管理・メール送信）",
        "LINE Official連携（お知らせ・アラート・前日自動確認・Webhook受付）",
        "自動受付エンジン（LINE/メール/FAX OCR→サジェスト）",
        "実績管理・日払い管理（手数料・月末自動差引）",
        "契約条件設定・各種CSV出力",
    ]
    add_rect(s6, 0.4, 1.3, 12.5, 0.45, GREEN)
    add_text_box(s6, "✓ Must — 6月20日リリース必須（35機能相当）",
                 0.5, 1.3, 12.3, 0.45, font_size=13, bold=True, color=WHITE)
    for i, item in enumerate(must_items):
        y = 1.9 + i * 0.56
        add_rect(s6, 0.4, y, 12.5, 0.48, RGBColor(0xDC, 0xFC, 0xE7))
        add_text_box(s6, f"✓  {item}", 0.55, y + 0.04, 12.2, 0.4,
                     font_size=11, color=DARK)

    # ====== スライド7: スケジュール ======
    s7 = prs.slides.add_slide(blank_layout)
    slide_header(s7, "05. スケジュール（34日計画）", "5月17日〜6月20日")
    week_data = [
        ("Week 1", "5/17-5/23", NAVY, "基盤構築", "DB設計・認証・マルチテナント・マスターデータ・権限設定"),
        ("Week 2", "5/24-5/30", RGBColor(0x05, 0x96, 0x69), "管制機能（前半）", "状況監視・自動監視・1号/2号業務配置・ホワイトボード"),
        ("Week 3", "5/31-6/06", RGBColor(0xD9, 0x77, 0x06), "管制機能（後半）＋隊員PWA", "月間シフト・管制表出力・業務レポート・隊員PWA4機能"),
        ("Week 4", "6/07-6/13", RGBColor(0xBE, 0x18, 0x5D), "LINE連携・自動受付・請求", "LINE Official・前日確認・メール/FAX自動受付・OCR・請求"),
        ("Week 5", "6/14-6/20", RGBColor(0x70, 0x1A, 0x75), "日払い・実績・テスト・デプロイ", "日払い管理・実績・CSV・結合テスト・UI調整・本番デプロイ"),
    ]
    for i, (week, period, color, phase, tasks_str) in enumerate(week_data):
        y = 1.3 + i * 1.18
        add_rect(s7, 0.4, y, 2.2, 1.0, color)
        add_text_box(s7, week, 0.45, y + 0.05, 2.1, 0.45,
                     font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(s7, period, 0.45, y + 0.5, 2.1, 0.4,
                     font_size=10, color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(s7, 2.7, y, 10.2, 1.0, RGBColor(0xF0, 0xF4, 0xFF))
        add_text_box(s7, phase, 2.85, y + 0.05, 9.9, 0.4,
                     font_size=13, bold=True, color=color)
        add_text_box(s7, tasks_str, 2.85, y + 0.5, 9.9, 0.4,
                     font_size=9, color=GRAY)

    # ====== スライド8: リスク管理 ======
    s8 = prs.slides.add_slide(blank_layout)
    slide_header(s8, "07. リスク管理", "主要リスクと対応策")
    risk_items = [
        ("🔴", "スコープ超過", "毎週進捗確認・Should→v2移動"),
        ("🔴", "LINE審査遅延", "今週中申請・モック実装で並行開発"),
        ("🟡", "特許侵害リスク", "独自設計（QRコード承認）で完全回避"),
        ("🔴", "給与計算不正確", "基本形態のみ・社労士確認後にv2対応"),
        ("🟡", "FAX設定遅延", "Week1中に契約完了・他チャネルで代替"),
        ("🟢", "API費用超過", "コスト上限設定・キャッシュ実装"),
    ]
    for i, (level, risk, action) in enumerate(risk_items):
        col_s = i % 2
        row_s = i // 2
        x = 0.4 + col_s * 6.5
        y = 1.4 + row_s * 1.9
        colors_map = {"🔴": (RGBColor(0xFE, 0xE2, 0xE2), RED),
                      "🟡": (RGBColor(0xFE, 0xF9, 0xC3), YELLOW),
                      "🟢": (RGBColor(0xDC, 0xFC, 0xE7), GREEN)}
        bg, fg = colors_map[level]
        add_rect(s8, x, y, 6.0, 1.7, bg)
        add_text_box(s8, f"{level} {risk}", x + 0.2, y + 0.1, 5.6, 0.55,
                     font_size=14, bold=True, color=DARK)
        add_text_box(s8, f"→ {action}", x + 0.2, y + 0.75, 5.6, 0.8,
                     font_size=10, color=GRAY)

    # ====== スライド9: 次のアクション ======
    s9 = prs.slides.add_slide(blank_layout)
    slide_header(s9, "08. 次のアクション", "今日・今週中にやること")
    actions = [
        ("🔴 今日中", [
            "LINE Official Account 開設申請",
            "Claude Code Maxプラン加入",
            "GitHubリポジトリ作成（guardsync）",
            "Railway新規プロジェクト作成",
        ]),
        ("🟡 今週中（5/23まで）", [
            "FAX→メール変換サービス契約（eFax等）",
            "LINE Messaging APIチャネル作成",
            "Google Vision API有効化・課金設定",
            "Week1タスク（基盤構築）完了",
        ]),
    ]
    for i, (period, items) in enumerate(actions):
        x = 0.4 + i * 6.5
        add_rect(s9, x, 1.3, 6.0, 0.5, NAVY if i == 0 else RGBColor(0xD9, 0x77, 0x06))
        add_text_box(s9, period, x + 0.1, 1.3, 5.8, 0.5,
                     font_size=14, bold=True, color=WHITE)
        for j, item in enumerate(items):
            y = 1.95 + j * 1.1
            add_rect(s9, x, y, 6.0, 0.9, RGBColor(0xF8, 0xFA, 0xFF))
            add_rect(s9, x, y, 0.15, 0.9,
                     NAVY if i == 0 else RGBColor(0xD9, 0x77, 0x06))
            add_text_box(s9, f"☐  {item}", x + 0.25, y + 0.1, 5.6, 0.65,
                         font_size=11, color=DARK)

    # ====== スライド10: Q&A ======
    s10 = prs.slides.add_slide(blank_layout)
    add_rect(s10, 0, 0, 13.33, 7.5, NAVY)
    add_text_box(s10, "ご確認・ご承認をお願いします", 1, 2.5, 11, 1.0,
                 font_size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(s10, "GuardSync プロジェクト計画書 v1.0",
                 1, 3.7, 11, 0.6, font_size=16, color=LIGHT_BLUE,
                 align=PP_ALIGN.CENTER)
    add_text_box(s10, "MINリリース：2026年6月20日　|　MAXリリース：2026年8月末",
                 1, 4.4, 11, 0.5, font_size=13, color=LIGHT_BLUE,
                 align=PP_ALIGN.CENTER)

    path = os.path.join(OUTPUT_DIR, "20260518_GuardSync_プロジェクト計画書.pptx")
    prs.save(path)
    print(f"PowerPoint出力完了: {path}")
    return path


# ============================================================
# メイン実行
# ============================================================
if __name__ == "__main__":
    print("GuardSync プロジェクト計画書を生成中...")
    print()
    word_path = create_word()
    excel_path = create_excel()
    pptx_path = create_pptx()
    print()
    print("=" * 60)
    print("✓ 全ファイルの生成が完了しました")
    print(f"  Word:       {os.path.basename(word_path)}")
    print(f"  Excel:      {os.path.basename(excel_path)}")
    print(f"  PowerPoint: {os.path.basename(pptx_path)}")
    print(f"  保存先: {OUTPUT_DIR}")
