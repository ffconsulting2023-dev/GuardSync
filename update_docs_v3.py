#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
GuardSync プロジェクト計画書 v1.2 更新スクリプト
LINE Official → LINE Works への変更を反映
"""

import os

OUTPUT_DIR = "/Users/t.take/AI management/project/output/GuardSync"


def update_word():
    from docx import Document
    from docx.shared import Pt, RGBColor, Cm
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    path = os.path.join(OUTPUT_DIR, "20260518_GuardSync_プロジェクト計画書.docx")
    doc = Document(path)

    def add_heading(text, level=1, color=None):
        h = doc.add_heading(text, level=level)
        if color:
            for run in h.runs:
                run.font.color.rgb = RGBColor(*color)
        return h

    def add_table(headers, rows):
        table = doc.add_table(rows=1, cols=len(headers))
        table.style = 'Table Grid'
        hdr = table.rows[0].cells
        for i, h in enumerate(headers):
            hdr[i].text = h
            hdr[i].paragraphs[0].runs[0].font.bold = True
            hdr[i].paragraphs[0].runs[0].font.size = Pt(9)
            hdr[i].paragraphs[0].alignment = WD_ALIGN_PARAGRAPH.CENTER
        for row_data in rows:
            row = table.add_row().cells
            for i, val in enumerate(row_data):
                row[i].text = str(val)
                row[i].paragraphs[0].runs[0].font.size = Pt(9)
        return table

    doc.add_page_break()

    add_heading("付録：v1.2変更履歴 — LINE Works採用", 1, (0x1A, 0x56, 0xDB))

    p = doc.add_paragraph()
    p.add_run("変更日：2026年5月18日　|　変更内容：警備員とのやりとりをLINE WorksベースのBot APIに変更").font.bold = True

    add_heading("変更理由", 2)
    doc.add_paragraph(
        "警備員との業務コミュニケーションにLINE Worksを採用する。"
        "LINE Worksはビジネス向けの法人向けサービスであり、管理者による隊員アカウントの一元管理、"
        "強力な既読確認機能、組織・チャンネル管理、企業向けセキュリティを標準提供する。"
        "個人のLINEアカウントと業務を分離できる点も隊員のプライバシー保護に有効である。"
    )

    add_heading("LINE vs LINE Works 比較", 2)
    add_table(
        ["項目", "LINE Official（旧）", "LINE Works（新）"],
        [
            ["用途", "一般消費者向け", "ビジネス・法人向け"],
            ["API", "Messaging API", "LINE Works Bot API v2"],
            ["アカウント管理", "個人アカウント", "管理者が全員を一元管理"],
            ["既読確認", "困難", "完全対応"],
            ["グループ管理", "限定的", "チャンネル・組織管理"],
            ["チャット", "独自実装が必要", "ネイティブ機能を利用（開発不要）"],
            ["認証", "LINE Login", "LINE Works OAuth 2.0"],
            ["セキュリティ", "標準", "企業向け強化"],
            ["料金", "無料〜", "フリー〜¥960/人/月（100名まで無料）"],
        ]
    )

    doc.add_paragraph()
    add_heading("LINE Works連携 機能構成", 2)
    add_table(
        ["機能", "実装方式", "備考"],
        [
            ["お知らせ配信", "LINE Works Bot API", "既読確認・返信受付が標準対応"],
            ["アラート通知（遅刻・未報告）", "LINE Works Bot API + Cron", "条件発火で自動送信"],
            ["前日自動確認送信", "LINE Works Bot API + Cron", "毎朝10時に翌日配置を自動送信"],
            ["自動受付Webhook", "LINE Works Webhook", "メッセージ受信→サジェスト"],
            ["グループチャット（現場別）", "LINE Worksネイティブ", "開発不要・管理者がチャンネル作成"],
            ["1対1チャット（管制↔隊員）", "LINE Worksネイティブ", "開発不要"],
            ["日払い申請", "LINE Works Bot フロー", "Botガイダンスで申請→システム記録"],
            ["休暇申請", "LINE Works Bot フロー", "同上"],
            ["出動・上下番連絡", "PWA（メイン）+ Bot（補助）", "PWAが主・Botは通知補助"],
        ]
    )

    doc.add_paragraph()
    add_heading("事前タスク 更新", 2)
    add_table(
        ["旧タスク", "新タスク", "期限"],
        [
            ["LINE Official Account 開設", "LINE Works テナント登録", "2026/05/18"],
            ["Messaging API チャネル作成", "LINE Works Bot 作成（Developer Console）", "2026/05/18"],
            ["LINE Login チャネル作成", "LINE Works OAuth 2.0 設定", "2026/05/18"],
        ]
    )

    doc.add_paragraph()
    add_heading("技術スタック 更新", 2)
    add_table(
        ["変更項目", "旧", "新"],
        [
            ["メッセージングAPI", "LINE Messaging API", "LINE Works Bot API v2"],
            ["認証", "LINE Login", "LINE Works OAuth 2.0"],
            ["Webhook", "LINE Webhook", "LINE Works Webhook"],
            ["チャット実装", "独自実装（WebSocket）", "LINE Worksネイティブ（実装不要）"],
        ]
    )

    path_save = os.path.join(OUTPUT_DIR, "20260518_GuardSync_プロジェクト計画書.docx")
    doc.save(path_save)
    print(f"Word更新完了: {path_save}")


def update_excel():
    from openpyxl import load_workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter

    path = os.path.join(OUTPUT_DIR, "20260518_GuardSync_プロジェクト計画書.xlsx")
    wb = load_workbook(path)

    NAVY = "1A56DB"
    GREEN = "05966940"
    LIGHT_GREEN = "DCFCE7"
    LIGHT_BLUE = "DBEAFE"
    WHITE = "FFFFFF"
    ORANGE = "EA580C"
    LIGHT_ORANGE = "FFEDD5"

    header_font = Font(bold=True, color="FFFFFF", size=10)
    header_fill = PatternFill("solid", fgColor="05966940"[0:6] if len("05966940") > 6 else "059669")
    body_font = Font(size=9)
    center_align = Alignment(horizontal="center", vertical="center", wrap_text=True)
    wrap_align = Alignment(vertical="center", wrap_text=True)
    thin_border = Border(
        left=Side(style="thin", color="CCCCCC"),
        right=Side(style="thin", color="CCCCCC"),
        top=Side(style="thin", color="CCCCCC"),
        bottom=Side(style="thin", color="CCCCCC"),
    )
    navy_fill = PatternFill("solid", fgColor=NAVY)
    green_fill = PatternFill("solid", fgColor="059669")

    # ====== 新シート: LINE Works連携設計 ======
    ws = wb.create_sheet("LINE Works連携設計")

    # タイトル
    ws.merge_cells("A1:F1")
    ws["A1"] = "LINE Works連携設計（v1.2 / 2026-05-18変更）"
    ws["A1"].font = Font(bold=True, size=14, color="FFFFFF")
    ws["A1"].fill = PatternFill("solid", fgColor="059669")
    ws["A1"].alignment = Alignment(horizontal="center", vertical="center")
    ws.row_dimensions[1].height = 35

    # 比較表
    ws.merge_cells("A2:F2")
    ws["A2"] = "■ LINE vs LINE Works 比較"
    ws["A2"].font = Font(bold=True, size=11, color="FFFFFF")
    ws["A2"].fill = PatternFill("solid", fgColor=NAVY)
    ws["A2"].alignment = Alignment(horizontal="center")

    comp_headers = ["項目", "LINE Official（旧）", "LINE Works（新）", "採用理由"]
    for col, h in enumerate(comp_headers, 1):
        c = ws.cell(row=3, column=col, value=h)
        c.font = Font(bold=True, color="FFFFFF", size=9)
        c.fill = PatternFill("solid", fgColor=NAVY)
        c.alignment = center_align
        c.border = thin_border

    comp_data = [
        ("用途", "一般消費者向け", "ビジネス・法人向け", "業務専用ツールとして分離できる"),
        ("API", "Messaging API", "LINE Works Bot API v2", "Bot APIで業務自動化"),
        ("アカウント管理", "個人アカウント", "管理者が全員を一元管理", "隊員の入退社管理が容易"),
        ("既読確認", "困難", "完全対応（標準機能）", "業務連絡の既読管理が必須"),
        ("グループ管理", "限定的", "チャンネル・組織管理", "現場別グループを管理者が作成"),
        ("チャット", "独自実装が必要", "ネイティブ機能を利用", "開発工数を大幅削減"),
        ("認証", "LINE Login", "LINE Works OAuth 2.0", "法人認証で信頼性向上"),
        ("セキュリティ", "標準", "企業向け強化", "業務データの保護"),
        ("料金", "無料〜", "フリー〜¥960/人/月", "100名まで無料・コスト効率良"),
    ]

    for i, row_data in enumerate(comp_data, start=4):
        fill_color = LIGHT_BLUE if i % 2 == 0 else WHITE
        for col, val in enumerate(row_data, 1):
            c = ws.cell(row=i, column=col, value=val)
            c.font = body_font
            c.fill = PatternFill("solid", fgColor=fill_color)
            c.border = thin_border
            c.alignment = wrap_align
            if col == 3:
                c.font = Font(size=9, bold=True, color="059669")
        ws.row_dimensions[i].height = 28

    # 機能連携設計
    ws.merge_cells("A14:F14")
    ws["A14"] = "■ LINE Works連携 機能設計"
    ws["A14"].font = Font(bold=True, size=11, color="FFFFFF")
    ws["A14"].fill = PatternFill("solid", fgColor="059669")
    ws["A14"].alignment = Alignment(horizontal="center")

    func_headers = ["機能", "実装方式", "トリガー", "送信先", "既読確認", "備考"]
    for col, h in enumerate(func_headers, 1):
        c = ws.cell(row=15, column=col, value=h)
        c.font = Font(bold=True, color="FFFFFF", size=9)
        c.fill = PatternFill("solid", fgColor=NAVY)
        c.alignment = center_align
        c.border = thin_border

    func_data = [
        ("お知らせ配信", "Bot API", "管制員が手動送信", "個人 or グループ", "✓", "返信も受付"),
        ("アラート通知（遅刻）", "Bot API + Cron", "時刻条件発火", "管制員・当該隊員", "✓", "未報告時に自動送信"),
        ("前日自動確認送信", "Bot API + Cron", "毎日10:00", "翌日出動隊員全員", "✓", "差別化機能"),
        ("LINE受付→サジェスト", "Webhook + AI", "メッセージ受信", "管制画面に表示", "-", "受注依頼の自動解析"),
        ("日払い申請", "Bot フロー", "隊員がメッセージ送信", "管制員に通知", "✓", "Botガイダンスで金額確認"),
        ("休暇申請", "Bot フロー", "隊員がメッセージ送信", "管制員に通知", "✓", "承認状況をBotで返信"),
        ("出動・上下番連絡", "PWA（メイン）", "隊員がボタンタップ", "管制側に自動反映", "-", "PWAが主・Botは補助通知"),
        ("グループチャット", "ネイティブ", "手動", "現場グループ", "✓", "開発不要・管理者がCh作成"),
        ("1対1チャット", "ネイティブ", "手動", "管制員↔隊員", "✓", "開発不要"),
    ]

    for i, row_data in enumerate(func_data, start=16):
        fill_color = LIGHT_GREEN if i % 2 == 0 else WHITE
        for col, val in enumerate(row_data, 1):
            c = ws.cell(row=i, column=col, value=val)
            c.font = body_font
            c.fill = PatternFill("solid", fgColor=fill_color)
            c.border = thin_border
            c.alignment = center_align if col in [2, 5] else wrap_align
        ws.row_dimensions[i].height = 28

    # 事前タスク更新
    ws.merge_cells("A26:F26")
    ws["A26"] = "■ 事前タスク 更新（LINE関連）"
    ws["A26"].font = Font(bold=True, size=11, color="FFFFFF")
    ws["A26"].fill = PatternFill("solid", fgColor=NAVY)
    ws["A26"].alignment = Alignment(horizontal="center")

    task_headers = ["旧タスク", "新タスク", "期限", "URL/備考"]
    for col, h in enumerate(task_headers, 1):
        c = ws.cell(row=27, column=col, value=h)
        c.font = Font(bold=True, color="FFFFFF", size=9)
        c.fill = PatternFill("solid", fgColor=NAVY)
        c.alignment = center_align
        c.border = thin_border

    task_updates = [
        ("LINE Official Account 開設", "LINE Works テナント登録", "2026/05/18", "works.line.biz で無料登録"),
        ("Messaging API チャネル作成", "LINE Works Bot 作成", "2026/05/18", "developers.worksmobile.com"),
        ("LINE Login チャネル作成", "LINE Works OAuth 2.0 設定", "2026/05/18", "Bot APIと同時に設定"),
    ]
    for i, row_data in enumerate(task_updates, start=28):
        for col, val in enumerate(row_data, 1):
            c = ws.cell(row=i, column=col, value=val)
            c.font = body_font
            c.fill = PatternFill("solid", fgColor=LIGHT_ORANGE)
            c.border = thin_border
            c.alignment = wrap_align
        ws.row_dimensions[i].height = 30

    # 列幅
    col_widths = [30, 25, 15, 25, 12, 25]
    for col, w in enumerate(col_widths, 1):
        ws.column_dimensions[get_column_letter(col)].width = w

    # ToDo シートのLINE関連タスクを更新
    ws_todo = wb["ToDo・事前タスク"]
    # 既存のLINE関連行を探して更新
    for row in ws_todo.iter_rows():
        for cell in row:
            if cell.value and "LINE Official" in str(cell.value):
                cell.value = cell.value.replace("LINE Official Account 開設申請", "LINE Works テナント登録（works.line.biz）")
            if cell.value and "Messaging API" in str(cell.value):
                cell.value = cell.value.replace("LINE Messaging APIチャネル作成", "LINE Works Bot作成（developers.worksmobile.com）")

    # 技術スタックシートの更新（概要シート）
    ws_overview = wb["概要"]
    for row in ws_overview.iter_rows():
        for cell in row:
            if cell.value and "LINE Official" in str(cell.value):
                cell.value = str(cell.value).replace("LINE Official Account + Messaging API", "LINE Works Bot API v2")

    wb.save(path)
    print(f"Excel更新完了: {path}")


def update_pptx():
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    path = os.path.join(OUTPUT_DIR, "20260518_GuardSync_プロジェクト計画書.pptx")
    prs = Presentation(path)

    blank_layout = prs.slide_layouts[6]

    NAVY = RGBColor(0x1A, 0x56, 0xDB)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT_BLUE = RGBColor(0xDB, 0xEA, 0xFE)
    DARK = RGBColor(0x1F, 0x29, 0x37)
    GRAY = RGBColor(0x6B, 0x72, 0x80)
    GREEN = RGBColor(0x05, 0x96, 0x69)
    LIGHT_GREEN = RGBColor(0xDC, 0xFC, 0xE7)
    ORANGE = RGBColor(0xEA, 0x58, 0x0C)
    LIGHT_ORANGE = RGBColor(0xFF, 0xED, 0xD5)

    def add_rect(slide, x, y, w, h, fill_color, line_color=None):
        shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        if line_color:
            shape.line.color.rgb = line_color
        else:
            shape.line.fill.background()
        return shape

    def add_text_box(slide, text, x, y, w, h, font_size=12, bold=False,
                     color=None, align=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color
        return txBox

    def slide_header(slide, title, subtitle=None, color=None):
        hdr_color = color or NAVY
        add_rect(slide, 0, 0, 13.33, 1.2, hdr_color)
        add_text_box(slide, title, 0.4, 0.15, 12, 0.7,
                     font_size=28, bold=True, color=WHITE)
        if subtitle:
            add_text_box(slide, subtitle, 0.4, 0.75, 12, 0.4,
                         font_size=13, color=LIGHT_BLUE)
        add_rect(slide, 0, 7.3, 13.33, 0.2, hdr_color)
        add_text_box(slide, "GuardSync プロジェクト計画書 v1.2 | 2026.05.18",
                     0.4, 7.3, 12, 0.2, font_size=8, color=WHITE,
                     align=PP_ALIGN.CENTER)

    # ====== スライド14: LINE Works採用 ======
    s14 = prs.slides.add_slide(blank_layout)
    slide_header(s14, "LINE Works 採用（v1.2変更）", "警備員とのコミュニケーション基盤", GREEN)

    # 変更バナー
    add_rect(s14, 0.3, 1.3, 12.7, 0.55, ORANGE)
    add_text_box(s14, "📢  LINE Official → LINE Works に変更  |  警備員との業務連絡をビジネス専用ツールに統一",
                 0.5, 1.35, 12.3, 0.45, font_size=12, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER)

    # 左: LINE Works のメリット
    add_rect(s14, 0.3, 2.0, 6.0, 5.2, LIGHT_GREEN)
    add_rect(s14, 0.3, 2.0, 6.0, 0.5, GREEN)
    add_text_box(s14, "✓  LINE Works を選ぶ理由", 0.45, 2.05, 5.7, 0.4,
                 font_size=14, bold=True, color=WHITE)

    merits = [
        ("👥", "管理者による隊員アカウント一元管理", "入退社時の即時アカウント停止が可能"),
        ("✅", "既読確認が完全対応（標準機能）", "業務連絡の既読を確実に把握"),
        ("💬", "チャット機能はネイティブで提供", "独自実装不要→開発工数を節約"),
        ("🔒", "企業向けセキュリティ強化", "業務データとプライベートを完全分離"),
        ("💰", "100名まで無料フリープラン", "初期コストゼロで始められる"),
    ]
    for j, (icon, title, sub) in enumerate(merits):
        y = 2.6 + j * 0.88
        add_rect(s14, 0.4, y, 5.7, 0.75, WHITE)
        add_text_box(s14, icon, 0.5, y + 0.1, 0.5, 0.55, font_size=18)
        add_text_box(s14, title, 1.05, y + 0.05, 4.9, 0.35,
                     font_size=11, bold=True, color=DARK)
        add_text_box(s14, sub, 1.05, y + 0.42, 4.9, 0.28,
                     font_size=9, color=GRAY)

    # 右: 連携機能一覧
    add_rect(s14, 6.8, 2.0, 6.2, 5.2, LIGHT_BLUE)
    add_rect(s14, 6.8, 2.0, 6.2, 0.5, NAVY)
    add_text_box(s14, "🔗  GuardSync × LINE Works 連携", 6.95, 2.05, 5.9, 0.4,
                 font_size=14, bold=True, color=WHITE)

    integrations = [
        ("Bot API", "お知らせ配信・アラート通知・前日自動確認"),
        ("Bot API", "日払い申請・休暇申請（Botフロー）"),
        ("Webhook", "受注依頼の自動受付→サジェスト"),
        ("ネイティブ", "グループチャット（現場別チャンネル）"),
        ("ネイティブ", "1対1チャット（管制↔隊員）"),
        ("OAuth 2.0", "隊員の認証・アカウント管理"),
    ]
    for j, (method, desc) in enumerate(integrations):
        y = 2.6 + j * 0.78
        add_rect(s14, 6.9, y, 5.9, 0.65, WHITE)
        add_rect(s14, 6.9, y, 1.2, 0.65,
                 GREEN if "ネイティブ" in method else NAVY)
        add_text_box(s14, method, 6.92, y + 0.1, 1.15, 0.45,
                     font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(s14, desc, 8.2, y + 0.1, 4.5, 0.45,
                     font_size=10, color=DARK)

    prs.save(path)
    print(f"PowerPoint更新完了: {path}")


if __name__ == "__main__":
    print("GuardSync プロジェクト計画書 v1.2 更新中（LINE Works変更）...")
    print()
    update_word()
    update_excel()
    update_pptx()
    print()
    print("=" * 60)
    print("✓ 全ファイルの更新が完了しました（v1.2）")
    print(f"  保存先: {OUTPUT_DIR}")
